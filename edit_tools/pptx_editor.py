"""Edit existing PPTX files — modify slides, delete, reorder."""

from pptx import Presentation
from copy import deepcopy


def edit_pptx_slide_text(file_path: str, slide_index: int, old_text: str, new_text: str, output_path: str | None = None) -> str:
    """Replace text in a specific slide."""
    prs = Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})")

    slide = prs.slides[slide_index]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        count += 1

    save_path = output_path or file_path
    prs.save(save_path)
    return f"Replaced {count} occurrence(s) in slide {slide_index}. Saved to {save_path}"


def delete_pptx_slide(file_path: str, slide_index: int, output_path: str | None = None) -> str:
    """Delete a slide at the given index."""
    prs = Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})")

    rId = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]

    save_path = output_path or file_path
    prs.save(save_path)
    return f"Slide {slide_index} deleted. Saved to {save_path}"


def reorder_pptx_slides(file_path: str, new_order: list[int], output_path: str | None = None) -> str:
    """Reorder slides by providing new index order. E.g. [2, 0, 1] moves slide 2 to first."""
    prs = Presentation(file_path)
    slide_count = len(prs.slides)

    if sorted(new_order) != list(range(slide_count)):
        raise ValueError(f"new_order must be a permutation of [0..{slide_count-1}], got {new_order}")

    # Reorder by manipulating the sldIdLst
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    for item in items:
        sldIdLst.remove(item)
    for idx in new_order:
        sldIdLst.append(items[idx])

    save_path = output_path or file_path
    prs.save(save_path)
    return f"Slides reordered to {new_order}. Saved to {save_path}"
