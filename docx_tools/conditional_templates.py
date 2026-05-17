"""Conditional templates for DOCX and PPTX — if/else blocks and loops."""

import re
from docx import Document
from pptx import Presentation


def render_docx_template(file_path: str, context: dict, output_path: str | None = None) -> str:
    """Render a DOCX template with conditional blocks and loops.

    Supports:
    - {{variable}} — simple replacement
    - {{#if condition}}...{{/if}} — conditional block (truthy check)
    - {{#unless condition}}...{{/unless}} — inverse conditional
    - {{#each items}}...{{/each}} — loop over list (use {{.}} for item)
    """
    doc = Document(file_path)

    for para in doc.paragraphs:
        para.text = _render_text(para.text, context)

    # Also process tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    para.text = _render_text(para.text, context)

    save_path = output_path or file_path
    doc.save(save_path)
    return f"Template rendered with {len(context)} variables. Saved to {save_path}"


def render_pptx_template(file_path: str, context: dict, output_path: str | None = None) -> str:
    """Render a PPTX template with conditional blocks and loops."""
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    rendered = _render_text(full_text, context)
                    if rendered != full_text:
                        # Clear runs and set new text
                        for run in para.runs:
                            run.text = ""
                        if para.runs:
                            para.runs[0].text = rendered

    save_path = output_path or file_path
    prs.save(save_path)
    return f"PPTX template rendered with {len(context)} variables. Saved to {save_path}"


def _render_text(text: str, context: dict) -> str:
    """Process template syntax in text."""
    # Process {{#each items}}...{{/each}}
    each_pattern = r"\{\{#each\s+(\w+)\}\}(.*?)\{\{/each\}\}"
    def replace_each(match):
        key = match.group(1)
        body = match.group(2)
        items = context.get(key, [])
        if not isinstance(items, list):
            return ""
        result = []
        for item in items:
            if isinstance(item, dict):
                rendered = body
                for k, v in item.items():
                    rendered = rendered.replace(f"{{{{{k}}}}}", str(v))
                result.append(rendered)
            else:
                result.append(body.replace("{{.}}", str(item)))
        return "\n".join(result)

    text = re.sub(each_pattern, replace_each, text, flags=re.DOTALL)

    # Process {{#if condition}}...{{/if}}
    if_pattern = r"\{\{#if\s+(\w+)\}\}(.*?)\{\{/if\}\}"
    def replace_if(match):
        key = match.group(1)
        body = match.group(2)
        value = context.get(key)
        if value:
            return body
        return ""

    text = re.sub(if_pattern, replace_if, text, flags=re.DOTALL)

    # Process {{#unless condition}}...{{/unless}}
    unless_pattern = r"\{\{#unless\s+(\w+)\}\}(.*?)\{\{/unless\}\}"
    def replace_unless(match):
        key = match.group(1)
        body = match.group(2)
        value = context.get(key)
        if not value:
            return body
        return ""

    text = re.sub(unless_pattern, replace_unless, text, flags=re.DOTALL)

    # Process simple {{variable}} replacements
    simple_pattern = r"\{\{(\w+)\}\}"
    def replace_simple(match):
        key = match.group(1)
        return str(context.get(key, match.group(0)))

    text = re.sub(simple_pattern, replace_simple, text)

    return text
