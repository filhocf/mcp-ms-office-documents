"""Merge multiple documents into one."""

from docx import Document
from pptx import Presentation
from copy import deepcopy


def merge_docx(file_paths: list[str], output_path: str) -> str:
    """Merge multiple DOCX files into one. Content is appended sequentially."""
    if not file_paths:
        raise ValueError("No files provided")

    merged = Document(file_paths[0])

    for path in file_paths[1:]:
        # Add page break between documents
        merged.add_page_break()
        sub_doc = Document(path)
        for element in sub_doc.element.body:
            merged.element.body.append(deepcopy(element))

    merged.save(output_path)
    return f"Merged {len(file_paths)} DOCX files into {output_path}"


def merge_pptx(file_paths: list[str], output_path: str) -> str:
    """Merge multiple PPTX files into one. Slides are appended sequentially."""
    if not file_paths:
        raise ValueError("No files provided")

    merged = Presentation(file_paths[0])

    for path in file_paths[1:]:
        src = Presentation(path)
        for slide in src.slides:
            layout = merged.slide_layouts[0]  # Use first layout as fallback
            new_slide = merged.slides.add_slide(layout)
            for shape in slide.shapes:
                el = deepcopy(shape.element)
                new_slide.shapes._spTree.append(el)

    merged.save(output_path)
    total_slides = len(merged.slides)
    return f"Merged {len(file_paths)} PPTX files ({total_slides} slides total) into {output_path}"
