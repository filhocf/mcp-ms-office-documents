"""Tests for conditional templates (DOCX + PPTX)."""

import pytest
from docx import Document
from pptx import Presentation

from docx_tools.conditional_templates import render_docx_template, render_pptx_template, _render_text


class TestRenderText:
    def test_simple_variable(self):
        assert _render_text("Hello {{name}}", {"name": "World"}) == "Hello World"

    def test_missing_variable_unchanged(self):
        assert _render_text("{{unknown}}", {}) == "{{unknown}}"

    def test_if_true(self):
        text = "{{#if show}}Visible{{/if}}"
        assert _render_text(text, {"show": True}) == "Visible"

    def test_if_false(self):
        text = "{{#if show}}Visible{{/if}}"
        assert _render_text(text, {"show": False}) == ""

    def test_unless_true(self):
        text = "{{#unless premium}}Free tier{{/unless}}"
        assert _render_text(text, {"premium": True}) == ""

    def test_unless_false(self):
        text = "{{#unless premium}}Free tier{{/unless}}"
        assert _render_text(text, {"premium": False}) == "Free tier"

    def test_each_simple(self):
        text = "{{#each items}}{{.}}, {{/each}}"
        result = _render_text(text, {"items": ["A", "B", "C"]})
        assert "A" in result and "B" in result and "C" in result

    def test_each_dict(self):
        text = "{{#each people}}{{name}} ({{age}}), {{/each}}"
        result = _render_text(text, {"people": [{"name": "Alice", "age": 30}, {"name": "Bob", "age": 25}]})
        assert "Alice (30)" in result
        assert "Bob (25)" in result

    def test_each_empty(self):
        text = "{{#each items}}{{.}}{{/each}}"
        assert _render_text(text, {"items": []}) == ""


class TestRenderDocxTemplate:
    def test_render_simple(self, tmp_path):
        path = str(tmp_path / "tpl.docx")
        doc = Document()
        doc.add_paragraph("Dear {{name}},")
        doc.add_paragraph("{{#if vip}}VIP access granted.{{/if}}")
        doc.save(path)

        out = str(tmp_path / "out.docx")
        render_docx_template(path, {"name": "Alice", "vip": True}, out)
        result = Document(out)
        assert result.paragraphs[0].text == "Dear Alice,"
        assert result.paragraphs[1].text == "VIP access granted."

    def test_render_loop(self, tmp_path):
        path = str(tmp_path / "tpl.docx")
        doc = Document()
        doc.add_paragraph("Items: {{#each items}}{{.}}, {{/each}}")
        doc.save(path)

        out = str(tmp_path / "out.docx")
        render_docx_template(path, {"items": ["X", "Y", "Z"]}, out)
        result = Document(out)
        assert "X" in result.paragraphs[0].text


class TestRenderPptxTemplate:
    def test_render_pptx(self, tmp_path):
        path = str(tmp_path / "tpl.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "{{title}}"
        slide.placeholders[1].text = "By {{author}}"
        prs.save(path)

        out = str(tmp_path / "out.pptx")
        render_pptx_template(path, {"title": "My Talk", "author": "Alice"}, out)
        result = Presentation(out)
        assert result.slides[0].shapes.title.text == "My Talk"
