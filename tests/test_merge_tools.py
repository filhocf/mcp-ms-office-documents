"""Tests for merge tools."""

import pytest
from docx import Document
from pptx import Presentation
from pathlib import Path

from merge_tools import merge_docx, merge_pptx


@pytest.fixture
def docx_files(tmp_path):
    paths = []
    for i in range(3):
        path = str(tmp_path / f"doc{i}.docx")
        doc = Document()
        doc.add_heading(f"Document {i}", level=1)
        doc.add_paragraph(f"Content of document {i}.")
        doc.save(path)
        paths.append(path)
    return paths


@pytest.fixture
def pptx_files(tmp_path):
    paths = []
    for i in range(3):
        path = str(tmp_path / f"pres{i}.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Presentation {i}"
        prs.save(path)
        paths.append(path)
    return paths


class TestMergeDocx:
    def test_merge_multiple(self, docx_files, tmp_path):
        out = str(tmp_path / "merged.docx")
        result = merge_docx(docx_files, out)
        assert "3" in result
        assert Path(out).exists()
        doc = Document(out)
        text = "\n".join(p.text for p in doc.paragraphs)
        assert "Document 0" in text
        assert "Document 1" in text
        assert "Document 2" in text

    def test_merge_single(self, docx_files, tmp_path):
        out = str(tmp_path / "single.docx")
        result = merge_docx([docx_files[0]], out)
        assert "1" in result

    def test_merge_empty_list(self):
        with pytest.raises(ValueError):
            merge_docx([], "/tmp/out.docx")


class TestMergePptx:
    def test_merge_multiple(self, pptx_files, tmp_path):
        out = str(tmp_path / "merged.pptx")
        result = merge_pptx(pptx_files, out)
        assert "3" in result
        prs = Presentation(out)
        assert len(prs.slides) >= 3

    def test_merge_single(self, pptx_files, tmp_path):
        out = str(tmp_path / "single.pptx")
        result = merge_pptx([pptx_files[0]], out)
        assert "1" in result

    def test_merge_empty_list(self):
        with pytest.raises(ValueError):
            merge_pptx([], "/tmp/out.pptx")
