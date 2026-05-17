"""Tests for edit_tools — editing existing DOCX, XLSX, PPTX files."""

import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from edit_tools import (
    edit_docx_paragraph, delete_docx_paragraph, search_replace_docx, insert_docx_paragraph,
    edit_xlsx_cell, insert_xlsx_row, delete_xlsx_row,
    edit_pptx_slide_text, delete_pptx_slide, reorder_pptx_slides,
)


@pytest.fixture
def sample_docx(tmp_path):
    path = str(tmp_path / "test.docx")
    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph("First paragraph.")
    doc.add_paragraph("Second paragraph with keyword here.")
    doc.save(path)
    return path


@pytest.fixture
def sample_xlsx(tmp_path):
    path = str(tmp_path / "test.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    ws.append(["Bob", 82])
    wb.save(path)
    return path


@pytest.fixture
def sample_pptx(tmp_path):
    path = str(tmp_path / "test.pptx")
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Slide One"
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])
    slide2.shapes.title.text = "Slide Two"
    slide3 = prs.slides.add_slide(prs.slide_layouts[0])
    slide3.shapes.title.text = "Slide Three"
    prs.save(path)
    return path


# === DOCX Edit Tests ===

class TestDocxEditor:
    def test_edit_paragraph(self, sample_docx, tmp_path):
        out = str(tmp_path / "out.docx")
        result = edit_docx_paragraph(sample_docx, 1, "Modified paragraph.", out)
        assert "updated" in result
        doc = Document(out)
        assert doc.paragraphs[1].text == "Modified paragraph."

    def test_edit_paragraph_preserves_style(self, sample_docx, tmp_path):
        out = str(tmp_path / "out.docx")
        edit_docx_paragraph(sample_docx, 0, "New Title", out)
        doc = Document(out)
        assert "Heading" in doc.paragraphs[0].style.name

    def test_delete_paragraph(self, sample_docx, tmp_path):
        out = str(tmp_path / "out.docx")
        doc_before = Document(sample_docx)
        count_before = len(doc_before.paragraphs)
        delete_docx_paragraph(sample_docx, 1, out)
        doc_after = Document(out)
        assert len(doc_after.paragraphs) == count_before - 1

    def test_search_replace(self, sample_docx, tmp_path):
        out = str(tmp_path / "out.docx")
        result = search_replace_docx(sample_docx, "keyword", "REPLACED", out)
        assert "1" in result  # 1 occurrence
        doc = Document(out)
        assert "REPLACED" in doc.paragraphs[2].text
        assert "keyword" not in doc.paragraphs[2].text

    def test_insert_paragraph(self, sample_docx, tmp_path):
        out = str(tmp_path / "out.docx")
        insert_docx_paragraph(sample_docx, 1, "Inserted text", "Normal", out)
        doc = Document(out)
        assert doc.paragraphs[1].text == "Inserted text"

    def test_edit_invalid_index(self, sample_docx):
        with pytest.raises(ValueError):
            edit_docx_paragraph(sample_docx, 999, "text")


# === XLSX Edit Tests ===

class TestXlsxEditor:
    def test_edit_cell(self, sample_xlsx, tmp_path):
        out = str(tmp_path / "out.xlsx")
        edit_xlsx_cell(sample_xlsx, "Data", "B2", "100", out)
        wb = load_workbook(out)
        assert wb["Data"]["B2"].value == "100"
        wb.close()

    def test_insert_row(self, sample_xlsx, tmp_path):
        out = str(tmp_path / "out.xlsx")
        insert_xlsx_row(sample_xlsx, "Data", 2, ["Charlie", "88"], out)
        wb = load_workbook(out)
        assert wb["Data"].cell(2, 1).value == "Charlie"
        assert wb["Data"].max_row == 4  # was 3, now 4
        wb.close()

    def test_delete_row(self, sample_xlsx, tmp_path):
        out = str(tmp_path / "out.xlsx")
        delete_xlsx_row(sample_xlsx, "Data", 2, out)
        wb = load_workbook(out)
        assert wb["Data"].max_row == 2  # was 3, now 2
        wb.close()

    def test_edit_invalid_sheet(self, sample_xlsx):
        with pytest.raises(ValueError):
            edit_xlsx_cell(sample_xlsx, "NonExistent", "A1", "x")


# === PPTX Edit Tests ===

class TestPptxEditor:
    def test_edit_slide_text(self, sample_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        result = edit_pptx_slide_text(sample_pptx, 0, "Slide One", "Modified Title", out)
        assert "1" in result
        prs = Presentation(out)
        assert prs.slides[0].shapes.title.text == "Modified Title"

    def test_delete_slide(self, sample_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        delete_pptx_slide(sample_pptx, 1, out)
        prs = Presentation(out)
        assert len(prs.slides) == 2

    def test_reorder_slides(self, sample_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        reorder_pptx_slides(sample_pptx, [2, 0, 1], out)
        prs = Presentation(out)
        assert prs.slides[0].shapes.title.text == "Slide Three"
        assert prs.slides[1].shapes.title.text == "Slide One"

    def test_reorder_invalid(self, sample_pptx):
        with pytest.raises(ValueError):
            reorder_pptx_slides(sample_pptx, [0, 1])  # missing index 2

    def test_delete_invalid_index(self, sample_pptx):
        with pytest.raises(ValueError):
            delete_pptx_slide(sample_pptx, 99)
